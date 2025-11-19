import os
from pypdf import PdfReader
from pptx import Presentation
from docx import Document

def extract_text(file_path: str, file_type: str) -> str:
    """
    Extract text from various file formats.
    """
    try:
        if file_type == "PDF":
            return _extract_from_pdf(file_path)
        elif file_type == "PPTX":
            return _extract_from_pptx(file_path)
        elif file_type == "DOCX":
            return _extract_from_docx(file_path)
        elif file_type == "TXT":
            return _extract_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    except Exception as e:
        print(f"Error parsing {file_path}: {str(e)}")
        return ""

def _extract_from_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def _extract_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def _extract_from_docx(file_path: str) -> str:
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def _extract_from_txt(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()
